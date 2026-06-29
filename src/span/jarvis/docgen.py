"""Documentgeneratie — Word/PowerPoint/Excel, optioneel uit een Lomans-template.

Een Office-TEMPLATE (.dotx/.potx/.xltx) heeft een ander content-type dan een
document; python-docx/pptx weigert 'm dan. `_template_to_document` zet het
content-type in [Content_Types].xml om zodat de template als bewerkbaar
document opengaat — zo behoudt de output de huisstijl (master/styles/logo).
"""

from __future__ import annotations

import io
import zipfile
from typing import Any

# template-content-type -> document-content-type, per soort
_CT_SWAP = {
    "word": ("wordprocessingml.template.main+xml", "wordprocessingml.document.main+xml"),
    "powerpoint": ("presentationml.template.main+xml", "presentationml.presentation.main+xml"),
    "excel": ("spreadsheetml.template.main+xml", "spreadsheetml.sheet.main+xml"),
}


def _template_to_document(raw: bytes, kind: str) -> bytes:
    src, dst = _CT_SWAP[kind]
    bin_in, bin_out = io.BytesIO(raw), io.BytesIO()
    with zipfile.ZipFile(bin_in) as zin, \
            zipfile.ZipFile(bin_out, "w", zipfile.ZIP_DEFLATED) as zout:
        for item in zin.infolist():
            data = zin.read(item.filename)
            if item.filename == "[Content_Types].xml":
                data = data.replace(src.encode(), dst.encode())
            zout.writestr(item, data)
    return bin_out.getvalue()


def _base(raw: bytes | None, kind: str):
    return io.BytesIO(_template_to_document(raw, kind)) if raw else None


# -- Word -----------------------------------------------------------------

def generate_docx(title: str, markdown: str, template_raw: bytes | None = None) -> bytes:
    from docx import Document
    doc = Document(_base(template_raw, "word")) if template_raw else Document()
    if title:
        try:
            doc.add_heading(title, level=0)
        except Exception:
            doc.add_paragraph(title)
    for line in (markdown or "").splitlines():
        s = line.rstrip()
        if not s:
            continue
        try:
            if s.startswith("### "):
                doc.add_heading(s[4:], level=3)
            elif s.startswith("## "):
                doc.add_heading(s[3:], level=2)
            elif s.startswith("# "):
                doc.add_heading(s[2:], level=1)
            elif s.lstrip().startswith(("- ", "* ")):
                doc.add_paragraph(s.lstrip()[2:], style="List Bullet")
            else:
                doc.add_paragraph(s)
        except Exception:
            doc.add_paragraph(s)
    out = io.BytesIO()
    doc.save(out)
    return out.getvalue()


# -- PowerPoint -----------------------------------------------------------

def parse_slides(markdown: str) -> list[dict[str, Any]]:
    """Markdown -> slides: elke '# ' (of '---') start een nieuwe slide; '- ' = bullet."""
    slides: list[dict[str, Any]] = []
    cur: dict[str, Any] | None = None
    for line in (markdown or "").splitlines():
        s = line.rstrip()
        if not s:
            continue
        if s == "---":
            cur = None
            continue
        if s.startswith("# ") or s.startswith("## "):
            cur = {"title": s.lstrip("# ").strip(), "bullets": []}
            slides.append(cur)
        else:
            if cur is None:
                cur = {"title": s[:80], "bullets": []}
                slides.append(cur)
            else:
                cur["bullets"].append(s.lstrip("-* ").strip())
    return slides


def generate_pptx(title: str, slides: list[dict[str, Any]],
                  template_raw: bytes | None = None) -> bytes:
    from pptx import Presentation
    prs = Presentation(_base(template_raw, "powerpoint")) if template_raw else Presentation()
    layouts = prs.slide_layouts
    title_layout = layouts[0]
    content_layout = layouts[1] if len(layouts) > 1 else layouts[0]
    # titeldia
    if title:
        s = prs.slides.add_slide(title_layout)
        if s.shapes.title:
            s.shapes.title.text = title
    for sl in slides:
        s = prs.slides.add_slide(content_layout)
        if s.shapes.title:
            s.shapes.title.text = sl.get("title", "")
        body = None
        for ph in s.placeholders:
            if ph.placeholder_format.idx != 0:  # niet de titel
                body = ph
                break
        if body is not None and sl.get("bullets"):
            tf = body.text_frame
            tf.clear()
            for i, b in enumerate(sl["bullets"]):
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                p.text = b
    out = io.BytesIO()
    prs.save(out)
    return out.getvalue()


# -- Excel ----------------------------------------------------------------

def parse_rows(content: str) -> list[list[str]]:
    """Tekst -> rijen: per regel gesplitst op tab, ';' of ',' (of '|' bij markdown-tabel)."""
    rows: list[list[str]] = []
    for line in (content or "").splitlines():
        s = line.strip()
        if not s or set(s) <= set("-|: "):  # markdown-scheidingsregel overslaan
            continue
        sep = "\t" if "\t" in s else ("|" if "|" in s else (";" if ";" in s else ","))
        cells = [c.strip() for c in s.strip("|").split(sep)]
        rows.append(cells)
    return rows


def generate_xlsx(title: str, rows: list[list[Any]], sheet_name: str = "Blad1",
                  template_raw: bytes | None = None) -> bytes:
    if template_raw:
        from openpyxl import load_workbook
        wb = load_workbook(_base(template_raw, "excel"))
        ws = wb.active
    else:
        from openpyxl import Workbook
        wb = Workbook()
        ws = wb.active
        ws.title = (sheet_name or "Blad1")[:31]
    for r in rows:
        ws.append(list(r))
    out = io.BytesIO()
    wb.save(out)
    return out.getvalue()
